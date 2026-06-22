#!/usr/bin/env python3
"""Genera la presentacion de la Entrega 3 de CampusLink (PPTX 16:9) con diseno pulido."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ---- Paleta de marca ----
NAVY      = RGBColor(0x1E, 0x3A, 0x5F)
NAVY2     = RGBColor(0x1E, 0x4F, 0x7A)
NAVYDK    = RGBColor(0x0F, 0x24, 0x40)
TEAL      = RGBColor(0x2E, 0xC4, 0xB6)
TEALDK    = RGBColor(0x1A, 0x9A, 0x8E)
LIGHTBLUE = RGBColor(0xA8, 0xC4, 0xE0)
SLATE     = RGBColor(0x64, 0x74, 0x8B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x1F, 0x2A, 0x37)
LIGHTBG   = RGBColor(0xF1, 0xF5, 0xF9)
CARDBG    = RGBColor(0xF8, 0xFA, 0xFC)

HEX = dict(NAVY="1E3A5F", NAVY2="1E4F7A", NAVYDK="0F2440", TEAL="2EC4B6",
           TEALDK="1A9A8E", LIGHTBG="F1F5F9", SLATE="64748B")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
DARK_SLIDES = []          # slides con fondo oscuro (para numero de pagina claro)
EXPO = {}                 # quien expone cada slide -> va al guion (notas), no a la vista
A = 'http://schemas.openxmlformats.org/drawingml/2006/main'


# ---------- helpers de bajo nivel (XML) ----------
def _grad_fill(fill_format, c1, c2, angle=90):
    xpr = fill_format._xPr
    for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill', 'a:grpFill'):
        e = xpr.find(qn(tag))
        if e is not None:
            xpr.remove(e)
    ang = int(angle * 60000)
    xml = (f'<a:gradFill xmlns:a="{A}" rotWithShape="1"><a:gsLst>'
           f'<a:gs pos="0"><a:srgbClr val="{c1}"/></a:gs>'
           f'<a:gs pos="100000"><a:srgbClr val="{c2}"/></a:gs>'
           f'</a:gsLst><a:lin ang="{ang}" scaled="1"/></a:gradFill>')
    el = etree.fromstring(xml)
    # el relleno debe ir antes de <a:ln> (formas) o de <a:effectLst> (fondo bgPr);
    # python-pptx deja un <a:effectLst/> en el bgPr, asi que no se puede solo append.
    ref = xpr.find(qn('a:ln'))
    if ref is None:
        ref = xpr.find(qn('a:effectLst'))
    if ref is not None:
        ref.addprevious(el)
    else:
        xpr.append(el)


def _shadow(shape, blur=10, dist=4, direction=5400000, alpha=74, color="0F2440"):
    spPr = shape._element.spPr
    ex = spPr.find(qn('a:effectLst'))
    if ex is not None:
        spPr.remove(ex)
    xml = (f'<a:effectLst xmlns:a="{A}"><a:outerShdw blurRad="{blur*12700}" '
           f'dist="{dist*12700}" dir="{direction}" rotWithShape="0">'
           f'<a:srgbClr val="{color}"><a:alpha val="{alpha*1000}"/></a:srgbClr>'
           f'</a:outerShdw></a:effectLst>')
    spPr.append(etree.fromstring(xml))


def _alpha_fill(shape, hexcolor, alpha):
    """Relleno solido translucido (para circulos decorativos)."""
    spPr = shape._element.spPr
    for tag in ('a:noFill', 'a:solidFill', 'a:gradFill'):
        e = spPr.find(qn(tag))
        if e is not None:
            spPr.remove(e)
    xml = (f'<a:solidFill xmlns:a="{A}"><a:srgbClr val="{hexcolor}">'
           f'<a:alpha val="{alpha*1000}"/></a:srgbClr></a:solidFill>')
    el = etree.fromstring(xml)
    ref = spPr.find(qn('a:ln'))
    if ref is None:
        ref = spPr.find(qn('a:effectLst'))
    if ref is not None:
        ref.addprevious(el)
    else:
        spPr.append(el)


# ---------- helpers de alto nivel ----------
def _box(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def _rect(slide, l, t, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if color is not None:
        sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    return sp


def _circle(slide, l, t, d, hexcolor, alpha):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    c.line.fill.background(); c.shadow.inherit = False
    _alpha_fill(c, hexcolor, alpha)
    return c


def _para(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT,
          space_after=6, first=False, level=0, spacing=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.level = level
    run = p.add_run(); run.text = text
    f = run.font
    f.size = Pt(size); f.bold = bold; f.color.rgb = color
    f.name = 'Calibri'
    if spacing is not None:   # letter spacing en 1/100 pt
        run._r.get_or_add_rPr().set('spc', str(int(spacing * 100)))
    return p


def notes(slide, text):
    expo = EXPO.get(id(slide))
    if expo:
        text = f"[Habla: {expo}] " + text
    slide.notes_slide.notes_text_frame.text = text


def _badge(slide, l, t, d, num):
    """Circulo teal con numero blanco."""
    c = _rect(slide, l, t, d, d, TEAL, MSO_SHAPE.OVAL)
    _shadow(c, blur=6, dist=2, alpha=55, color="1A9A8E")
    c.text_frame.word_wrap = False
    c.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _para(c.text_frame, str(num), 14, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    return c


def header(slide, kicker, title, expositor=None):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = WHITE
    band = _rect(slide, 0, 0, 13.333, 1.5, NAVY)
    _grad_fill(band.fill, HEX["NAVY"], HEX["NAVY2"], angle=0)
    _rect(slide, 0, 1.5, 13.333, 0.09, TEAL)
    # acento cuadrado teal junto al kicker
    _rect(slide, 0.6, 0.34, 0.16, 0.16, TEAL)
    tf = _box(slide, 0.88, 0.22, 11, 0.45).text_frame
    _para(tf, kicker.upper(), 12.5, TEAL, bold=True, first=True, spacing=2.2)
    tf2 = _box(slide, 0.58, 0.62, 12.1, 0.82).text_frame
    _para(tf2, title, 29, WHITE, bold=True, first=True)
    EXPO[id(slide)] = expositor
    return slide


def bullets(slide, items, left=0.75, top=1.95, width=11.9, size=17, gap=11):
    tf = _box(slide, left, top, width, 5.0).text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        txt, lvl = it if isinstance(it, tuple) else (it, 0)
        p = _para(tf, ("▸  " if lvl == 0 else "•  ") + txt,
                  size if lvl == 0 else size-2,
                  DARK if lvl == 0 else SLATE, bold=(lvl == 0),
                  first=(i == 0), space_after=gap, level=lvl)
        p.runs[0].font.color.rgb = TEAL if lvl == 0 else SLATE
        # solo la vinheta en teal; el texto en color normal
        # (separamos vinheta y texto en dos runs)
        full = p.runs[0].text
        p.runs[0].text = full[:3]
        r2 = p.add_run(); r2.text = full[3:]
        r2.font.size = Pt(size if lvl == 0 else size-2)
        r2.font.bold = (lvl == 0)
        r2.font.color.rgb = DARK if lvl == 0 else SLATE
        r2.font.name = 'Calibri'


def card(slide, l, t, w, h, fill=WHITE, accent=True):
    cd = _rect(slide, l, t, w, h, fill, MSO_SHAPE.ROUNDED_RECTANGLE)
    _shadow(cd, blur=11, dist=4, alpha=80)
    if accent:
        _rect(slide, l, t, w, 0.13, TEAL, MSO_SHAPE.ROUNDED_RECTANGLE)
    return cd


def section(title, num, expositor, subtitle=None):
    s = prs.slides.add_slide(BLANK); DARK_SLIDES.append(s)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
    _grad_fill(s.background.fill, HEX["NAVY"], HEX["NAVYDK"], angle=55)
    _circle(s, 9.2, -1.6, 5.6, HEX["TEAL"], 9)
    _circle(s, 10.6, 4.7, 3.4, HEX["TEAL"], 7)
    _circle(s, -1.3, 4.4, 3.6, HEX["NAVY2"], 40)
    # numero fantasma grande
    tfg = _box(s, 0.55, 0.7, 6, 2.6).text_frame
    _para(tfg, f"{num:02d}", 150, NAVY2, bold=True, first=True)
    _rect(s, 0.95, 3.5, 1.9, 0.14, TEAL)
    tf2 = _box(s, 0.9, 3.75, 11.5, 1.6).text_frame; tf2.word_wrap = True
    _para(tf2, title, 42, WHITE, bold=True, first=True)
    if subtitle:
        tf3 = _box(s, 0.95, 5.0, 10.8, 1.0).text_frame; tf3.word_wrap = True
        _para(tf3, subtitle, 18, LIGHTBLUE, first=True)
    notes(s, f"[Habla: {expositor}] {subtitle or title}. "
             "Hacer un pase natural desde quien venia hablando antes de entrar a esta seccion.")
    return s


def table_slide(slide, headers, rows, top=1.95, col_widths=None, font=13):
    nrows, ncols = len(rows) + 1, len(headers)
    left, width = Inches(0.75), Inches(11.83)
    height = Inches(min(4.9, 0.55 + 0.46 * nrows))
    gtbl = slide.shapes.add_table(nrows, ncols, left, Inches(top), width, height)
    tbl = gtbl.table
    tbl.first_row = False
    # quitar estilo por defecto (banding propio)
    tblPr = tbl._tbl.tblPr
    tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.2)
        p = c.text_frame.paragraphs[0]
        r = p.add_run(); r.text = h
        r.font.size = Pt(font + 1); r.font.bold = True; r.font.color.rgb = WHITE
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = LIGHTBG if i % 2 else WHITE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Inches(0.2)
            p = c.text_frame.paragraphs[0]
            rr = p.add_run(); rr.text = str(val)
            rr.font.size = Pt(font); rr.font.color.rgb = DARK
            rr.font.bold = (j == 1 and len(str(val)) <= 3)  # resalta numeros
            if rr.font.bold:
                rr.font.color.rgb = TEALDK
            rr.font.name = 'Calibri'
    return tbl


def stat_cards(slide, cards, top=4.25):
    n = len(cards); total_w = 11.83; gap = 0.32
    cw = (total_w - gap * (n - 1)) / n
    x = 0.75
    for num, lab in cards:
        cd = _rect(slide, x, top, cw, 1.95, NAVY, MSO_SHAPE.ROUNDED_RECTANGLE)
        _grad_fill(cd.fill, HEX["NAVY2"], HEX["NAVY"], angle=90)
        _shadow(cd, blur=11, dist=4, alpha=78)
        _rect(slide, x + cw/2 - 0.35, top + 0.28, 0.7, 0.08, TEAL)
        tf = _box(slide, x, top + 0.42, cw, 1.4).text_frame; tf.word_wrap = True
        _para(tf, num, 36, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=2)
        _para(tf, lab, 12.5, LIGHTBLUE, align=PP_ALIGN.CENTER)
        x += cw + gap


def demo_slide(num, title, expositor, narrativa, captura):
    s = prs.slides.add_slide(BLANK)
    header(s, f"Demostración · Paso {num}", title, expositor)
    # marco de captura
    ph = _rect(s, 0.75, 2.0, 6.7, 4.7, CARDBG, MSO_SHAPE.ROUNDED_RECTANGLE,
               line=LIGHTBLUE)
    _shadow(ph, blur=10, dist=3, alpha=70)
    tfp = ph.text_frame; tfp.word_wrap = True; tfp.vertical_anchor = MSO_ANCHOR.MIDDLE
    _para(tfp, "📸", 42, SLATE, align=PP_ALIGN.CENTER, first=True, space_after=6)
    _para(tfp, captura, 13, SLATE, align=PP_ALIGN.CENTER)
    # tarjeta de narrativa
    nc = card(s, 7.75, 2.0, 4.83, 4.7, fill=WHITE, accent=False)
    _rect(s, 7.75, 2.0, 0.13, 4.7, TEAL, MSO_SHAPE.ROUNDED_RECTANGLE)
    th = _box(s, 8.1, 2.2, 4.3, 0.4).text_frame
    _para(th, "NARRATIVA", 12.5, TEALDK, bold=True, first=True, spacing=2)
    for i, step in enumerate(narrativa):
        yb = 2.78 + i * 0.92
        _badge(s, 8.1, yb, 0.42, i + 1)
        tb = _box(s, 8.68, yb - 0.06, 3.75, 0.85).text_frame; tb.word_wrap = True
        _para(tb, step, 13.5, DARK, first=True)
    notes(s, f"DEMO EN VIVO — {title}. Narrar mientras se navega, explicando el porque de "
             "cada accion (no solo los clics). Pasos: " + " | ".join(narrativa) +
             ". Tener captura de respaldo por si falla la demo.")
    return s


def footer(slide, idx, dark):
    col = LIGHTBLUE if dark else SLATE
    tf = _box(slide, 0.6, 7.04, 8, 0.35).text_frame
    _para(tf, "CampusLink · Entrega 3 · INF331", 9, col, first=True)
    tf2 = _box(slide, 12.0, 7.04, 0.85, 0.35).text_frame
    _para(tf2, f"{idx:02d}", 10, col, bold=True, align=PP_ALIGN.RIGHT, first=True)


# ============================================================
# BLOQUE 1 — INTRODUCCIÓN (PORTADA)
# ============================================================
s = prs.slides.add_slide(BLANK); DARK_SLIDES.append(s)
s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
_grad_fill(s.background.fill, HEX["NAVY"], HEX["NAVYDK"], angle=60)
_circle(s, 8.7, -2.0, 6.8, HEX["TEAL"], 10)
_circle(s, 10.8, 4.4, 3.8, HEX["TEAL"], 8)
_circle(s, -1.8, 5.0, 4.0, HEX["NAVY2"], 45)
_rect(s, 0, 0, 0.32, 7.5, TEAL)
tf = _box(s, 1.0, 1.95, 11.5, 1.5).text_frame
_para(tf, "CampusLink", 58, WHITE, bold=True, first=True)
tf2 = _box(s, 1.0, 3.5, 11.5, 0.8).text_frame
_para(tf2, "Plataforma de Oportunidades Universitarias", 23, TEAL, first=True)
# tarjeta de info inferior
info = _rect(s, 1.0, 4.75, 8.6, 1.7, NAVYDK, MSO_SHAPE.ROUNDED_RECTANGLE)
_alpha_fill(info, "0F2440", 55); _shadow(info, blur=12, dist=4, alpha=60)
tf3 = _box(s, 1.35, 4.95, 8.0, 1.4).text_frame
_para(tf3, "INF331 – Pruebas de Software   ·   Entrega 3", 15, WHITE, bold=True,
      first=True, space_after=7)
_para(tf3, "Pruebas E2E/UI con Playwright integradas al pipeline CI/CD", 14, TEAL,
      space_after=10)
_para(tf3, "José Meza · Matías Barraza · Leonardo Chacón · Aarón Vargas", 13, LIGHTBLUE)
notes(s, "Buenas, somos el equipo de CampusLink del curso de Pruebas de Software. "
         "Hoy les presentamos la Entrega 3, donde sumamos pruebas de interfaz E2E con "
         "Playwright y las dejamos corriendo automaticamente dentro de nuestro pipeline. "
         "[Expositor: Matías]")

# ============================================================
# BLOQUE 2 — CONTEXTO Y PROBLEMA
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "Contexto y problema", "¿Qué resuelve CampusLink?", "José")
c1 = card(s, 0.75, 2.0, 5.85, 4.55)
tfa = _box(s, 1.05, 2.35, 5.3, 4.1).text_frame; tfa.word_wrap = True
_para(tfa, "EL PROBLEMA", 13.5, TEALDK, bold=True, first=True, space_after=10, spacing=1.5)
for t in ["Las oportunidades académicas están dispersas",
          "Correos, grupos de WhatsApp, murales físicos",
          "Se pierden oportunidades por falta de difusión",
          "Desconexión entre quien publica y quien postula"]:
    _para(tfa, "•  " + t, 15.5, DARK, space_after=11)
c2 = card(s, 6.95, 2.0, 5.63, 4.55)
tfb = _box(s, 7.25, 2.35, 5.1, 4.1).text_frame; tfb.word_wrap = True
_para(tfb, "LA SOLUCIÓN", 13.5, TEALDK, bold=True, first=True, space_after=10, spacing=1.5)
for t in ["Centraliza las oportunidades en una sola plataforma",
          "Publicar con detalle, requisitos y fecha límite",
          "Buscar y filtrar por tipo de oportunidad",
          "Postular y seguir el estado de cada postulación",
          "Gestionar perfil e historial"]:
    _para(tfb, "•  " + t, 15.5, DARK, space_after=10)
notes(s, "El problema que atacamos: las oportunidades académicas —tutorias, ayudantias, "
         "practicas— viven dispersas en correos y grupos de WhatsApp, y se pierden. "
         "CampusLink las centraliza: quien publica lo hace con requisitos y fecha limite, "
         "y el estudiante busca, filtra y postula desde un solo lugar.")

s = prs.slides.add_slide(BLANK)
header(s, "Contexto y problema", "¿Cómo nos organizamos?", "José")
bullets(s, [
    "Metodología iterativa por entregas:",
    ("E1: CRUD + autenticación + pruebas unitarias (Jest)", 1),
    ("E2: pipeline CI/CD + 2 nuevos requisitos funcionales", 1),
    ("E3: pruebas E2E/UI con Playwright integradas al pipeline", 1),
    "Jira — tablero Kanban con issues priorizados y estimados",
    "Git Flow — ramas por feature y Pull Requests hacia main",
    "Discord — comunicación del equipo y alertas del pipeline",
    "GitHub Actions — CI automático en cada push / PR",
], top=2.0, gap=9)
notes(s, "Seguimos la misma metodologia de las entregas anteriores: Kanban en Jira con "
         "historias estimadas, Git Flow con ramas por feature y PRs, y Discord para "
         "comunicarnos. La novedad transversal es que todo cambio pasa por CI en GitHub "
         "Actions. Vamos avanzando por entregas: hoy cerramos la E3 con E2E.")

# ============================================================
# BLOQUE 3 — FUNCIONALIDADES Y ARQUITECTURA
# ============================================================
section("Funcionalidades y arquitectura", 3, "Leonardo",
        "Qué construimos y sobre qué tecnologías corre")

s = prs.slides.add_slide(BLANK)
header(s, "Funcionalidades", "Requisitos funcionales y no funcionales", "Leonardo")
card(s, 0.75, 2.0, 5.85, 4.55)
tfa = _box(s, 1.05, 2.35, 5.3, 4.1).text_frame; tfa.word_wrap = True
_para(tfa, "FUNCIONALES", 13.5, TEALDK, bold=True, first=True, space_after=9, spacing=1.5)
for t in ["Registro e inicio de sesión con JWT",
          "CRUD completo de oportunidades",
          "Búsqueda y filtrado por tipo / texto",
          "Postulación e historial por usuario",
          "Control de acceso basado en roles (RBAC)",
          "Perfil de usuario con estadísticas"]:
    _para(tfa, "•  " + t, 14.5, DARK, space_after=8.5)
card(s, 6.95, 2.0, 5.63, 4.55)
tfb = _box(s, 7.25, 2.35, 5.1, 4.1).text_frame; tfb.word_wrap = True
_para(tfb, "NO FUNCIONALES", 13.5, TEALDK, bold=True, first=True, space_after=9, spacing=1.5)
for t in ["API REST con prefijo /api",
          "Autenticación stateless (JWT)",
          "Validación de entradas (class-validator)",
          "Diseño responsive (PC y celular)",
          "Deploy continuo (Vercel + Railway)",
          "CI automatizado con GitHub Actions"]:
    _para(tfb, "•  " + t, 14.5, DARK, space_after=8.5)
notes(s, "Estos son los requisitos que cubrimos. A la izquierda los funcionales, con el "
         "CRUD y el control de acceso por roles como nucleo. A la derecha los no "
         "funcionales: API REST stateless, validacion de entradas y, clave para este "
         "curso, CI automatizado y deploy continuo.")

s = prs.slides.add_slide(BLANK)
header(s, "Lo nuevo en E2 y E3", "Funcionalidades agregadas sobre el MVP", "Leonardo")
bullets(s, [
    "Roles ampliados: ESTUDIANTE, DOCENTE, INSTITUCION, ADMIN",
    "Módulo de administración: suspender usuarios, bloquear oportunidades",
    "Sistema de notificaciones (resultado, feedback, modificaciones)",
    "Guardar oportunidades (favoritos)",
    "Gestión de postulantes + finalización del proceso + feedback",
    "Chat entre postulante y publicador",
    "Estados de oportunidad y postulación con transiciones automáticas",
    "Búsqueda insensible a acentos",
], top=2.0, gap=8.5, size=16)
notes(s, "La Entrega 1 era un MVP. En E2 y E3 la app crecio bastante: anadimos roles y un "
         "modulo de administracion, notificaciones, favoritos, gestion de postulantes con "
         "feedback, chat, y una maquina de estados. Todo esto tambien amplio la base de pruebas.")

s = prs.slides.add_slide(BLANK)
header(s, "Arquitectura", "Flujo de la aplicación", "Leonardo")
table_slide(s,
    ["Acción", "Ruta Frontend", "Endpoint Backend"],
    [["Listar", "/opportunities", "GET /api/opportunities"],
     ["Ver detalle", "/opportunities/:id", "GET /api/opportunities/:id"],
     ["Crear", "/opportunities/new", "POST /api/opportunities"],
     ["Postular", "Botón en detalle", "POST /api/applications"],
     ["Mis postulaciones", "/profile", "GET /api/applications/mine"],
     ["Notificaciones", "/notifications", "GET /api/notifications"]],
    top=2.05, col_widths=[3.3, 4.0, 4.53], font=14)
notes(s, "El flujo tipico: el usuario se registra, ve el listado, filtra, entra al detalle "
         "y postula o crea una oportunidad. Cada accion del frontend mapea a un endpoint "
         "REST del backend; aqui ven esa correspondencia directa.")

s = prs.slides.add_slide(BLANK)
header(s, "Arquitectura", "Arquitectura general y stack", "Leonardo")
cols = [
    ("FRONTEND — Vercel", ["React 19 + Vite 8", "TypeScript + React Router v7",
                            "Axios → capa src/api/", "Jest + Testing Library"]),
    ("BACKEND — Railway", ["NestJS 11", "Controller → Service → Repository",
                           "TypeORM + class-validator", "JWT stateless · RBAC"]),
    ("DATOS E INFRA", ["PostgreSQL 15 (Railway)", "GitHub Actions (CI)",
                       "npm workspaces (monorepo)", "Deploy continuo"]),
]
x = 0.75
for titulo, items in cols:
    cd = _rect(s, x, 2.05, 3.85, 4.5, WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
    _shadow(cd, blur=11, dist=4, alpha=80)
    head = _rect(s, x, 2.05, 3.85, 0.78, NAVY, MSO_SHAPE.ROUNDED_RECTANGLE)
    _grad_fill(head.fill, HEX["NAVY2"], HEX["NAVY"], angle=0)
    head.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _para(head.text_frame, titulo, 13, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    tc = _box(s, x + 0.28, 3.05, 3.35, 3.4).text_frame; tc.word_wrap = True
    for i, it in enumerate(items):
        p = _para(tc, "▸  ", 13.5, TEAL, bold=True, first=(i == 0), space_after=11)
        r = p.add_run(); r.text = it
        r.font.size = Pt(13.5); r.font.color.rgb = DARK; r.font.name = 'Calibri'
    x += 4.15
notes(s, "La arquitectura tiene tres piezas: el frontend React en Vercel, el backend "
         "NestJS en Railway con el patron Controller-Service-Repository, y PostgreSQL. "
         "Todo es un monorepo con npm workspaces, lo que deja compartir tipos entre front "
         "y back, y CI en GitHub Actions.")

# ============================================================
# BLOQUE 4 — ESTRATEGIA DE PRUEBAS (FOCO)
# ============================================================
section("Estrategia de pruebas", 4, "Aarón",
        "El corazón de la Entrega 3: del unit test al E2E automatizado")

s = prs.slides.add_slide(BLANK)
header(s, "Estrategia de testing", "¿Qué probamos y cómo?", "Aarón")
card(s, 0.75, 1.95, 5.85, 2.05)
tfa = _box(s, 1.05, 2.25, 5.3, 1.7).text_frame; tfa.word_wrap = True
_para(tfa, "QUÉ PROBAMOS", 13, TEALDK, bold=True, first=True, space_after=7, spacing=1.5)
for t in ["Lógica de negocio de los servicios (backend)",
          "Componentes UI y capa HTTP/API (frontend)",
          "Flujos completos de usuario (E2E)"]:
    _para(tfa, "•  " + t, 13.5, DARK, space_after=5)
card(s, 6.95, 1.95, 5.63, 2.05)
tfb = _box(s, 7.25, 2.25, 5.1, 1.7).text_frame; tfb.word_wrap = True
_para(tfb, "TÉCNICAS", 13, TEALDK, bold=True, first=True, space_after=7, spacing=1.5)
for t in ["Caja negra · Partición de equivalencia",
          "Valores límite", "Mocking de dependencias"]:
    _para(tfb, "•  " + t, 13.5, DARK, space_after=5)
stat_cards(s, [("139", "Unitarias backend"), ("120", "Unitarias frontend"),
               ("17", "E2E Playwright"), ("276", "Total · 100% verde")], top=4.4)
notes(s, "Nuestra estrategia combina niveles. Probamos la logica de los servicios en el "
         "backend, los componentes y la capa de API en el frontend, y ahora los flujos "
         "completos con E2E. Aplicamos caja negra, particion de equivalencia, valores "
         "limite y mocking. En total 276 pruebas, todas en verde: crecimos desde las 192 "
         "de la primera entrega.")

s = prs.slides.add_slide(BLANK)
header(s, "Herramienta unitaria", "Jest — pruebas unitarias", "Aarón")
bullets(s, [
    "Test runner + assertions + mocks en un solo paquete",
    "Soporte TypeScript (ts-jest en backend, babel-jest en frontend)",
    "Entorno jsdom para simular el navegador en el frontend",
    "Repositorios TypeORM mockeados → servicios aislados",
    "Reportes de cobertura integrados (--coverage)",
    "Se ejecuta en CI en cada push (job Jest Unit Tests)",
], top=2.0, gap=12)
notes(s, "Para las unitarias usamos Jest, en back y front. Nos da runner, aserciones y "
         "mocks en un solo lugar. En el backend mockeamos los repositorios de TypeORM "
         "para aislar la logica de negocio. Todo corre en CI en cada push.")

s = prs.slides.add_slide(BLANK)
header(s, "Herramienta E2E · Entrega 3", "Playwright — pruebas de interfaz", "Aarón")
card(s, 0.75, 2.0, 5.85, 4.55)
tfa = _box(s, 1.05, 2.35, 5.3, 4.1).text_frame; tfa.word_wrap = True
_para(tfa, "QUÉ ES Y CÓMO LO USAMOS", 13, TEALDK, bold=True, first=True, space_after=9, spacing=1.3)
for t in ["Automatiza un navegador real (Chromium) headless",
          "Simula al usuario: clics, formularios, navegación",
          "Selectores estables con data-testid",
          "Auto-waiting → menos tests intermitentes",
          "BASE_URL por variable de entorno (no hardcode)",
          "Reporte HTML + trace + screenshots en fallos"]:
    _para(tfa, "•  " + t, 14, DARK, space_after=8.5)
cd = _rect(s, 6.95, 2.0, 5.63, 4.55, NAVY, MSO_SHAPE.ROUNDED_RECTANGLE)
_grad_fill(cd.fill, HEX["NAVY2"], HEX["NAVY"], angle=70); _shadow(cd, blur=12, dist=4, alpha=78)
tfb = _box(s, 7.3, 2.35, 5.0, 4.1).text_frame; tfb.word_wrap = True
_para(tfb, "¿POR QUÉ PLAYWRIGHT Y NO SELENIUM?", 13, TEAL, bold=True, first=True,
      space_after=9, spacing=1.3)
for t in ["Mismo ecosistema JS/TS del proyecto",
          "Headless nativo, ideal para CI",
          "Reportes y trazas integrados sin plugins",
          "Auto-waiting reduce la fragilidad",
          "Era una alternativa permitida por el enunciado"]:
    _para(tfb, "✓  " + t, 14, WHITE, space_after=9)
notes(s, "La estrella de esta entrega es Playwright. Maneja un navegador real headless y "
         "simula al usuario de punta a punta. El enunciado proponia Selenium por defecto, "
         "pero elegimos Playwright porque es del mismo ecosistema JS/TS, tiene headless "
         "nativo y trae reportes, trazas y screenshots integrados. Su auto-waiting nos "
         "evito muchos tests intermitentes.")

s = prs.slides.add_slide(BLANK)
header(s, "Cobertura E2E", "Los 17 flujos críticos automatizados", "Aarón")
table_slide(s,
    ["Flujo", "Casos", "Qué valida"],
    [["Autenticación", "3", "Login válido / inválido, registro"],
     ["Navegación / guards", "1", "Ruta protegida redirige a login"],
     ["Crear oportunidad", "2", "Creación y validación de campos"],
     ["Editar / Eliminar", "3", "Carga datos, guarda, elimina"],
     ["Listar / Detalle", "4", "Listado, detalle y badge de estado"],
     ["Búsqueda y filtro", "3", "Por texto, por tipo, limpiar"],
     ["Postular", "1", "Estudiante postula con éxito"]],
    top=2.05, col_widths=[3.6, 1.4, 6.83], font=14)
notes(s, "17 pruebas E2E que cubren los flujos criticos de usuario final: autenticacion, "
         "el CRUD completo, busqueda y filtros, los guards de navegacion y la postulacion. "
         "El enunciado pedia al menos 10, asi que estamos holgados y enfocados en lo que "
         "de verdad importa al usuario.")

s = prs.slides.add_slide(BLANK)
header(s, "Integración CI/CD", "El pipeline en GitHub Actions", "Aarón")
jobs = [("1", "Backend build", "Type-check (tsc)"),
        ("2", "Frontend build", "Type-check + build"),
        ("3", "Jest", "Unitarias + PostgreSQL"),
        ("4", "Playwright E2E", "Levanta app + 17 E2E"),
        ("5", "Discord", "Notifica éxito / fallo")]
x = 0.6; cw = 2.3
for i, (n, jt, jd) in enumerate(jobs):
    is_e2e = (i == 3)
    cd = _rect(s, x, 2.35, cw, 1.85, NAVY, MSO_SHAPE.ROUNDED_RECTANGLE)
    if is_e2e:
        _grad_fill(cd.fill, HEX["TEAL"], HEX["TEALDK"], angle=90)
    else:
        _grad_fill(cd.fill, HEX["NAVY2"], HEX["NAVY"], angle=90)
    _shadow(cd, blur=9, dist=3, alpha=72)
    bd = _rect(s, x + cw/2 - 0.28, 2.5, 0.56, 0.56, WHITE if not is_e2e else NAVY,
               MSO_SHAPE.OVAL)
    bd.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _para(bd.text_frame, n, 15, NAVY if not is_e2e else WHITE, bold=True,
          align=PP_ALIGN.CENTER, first=True)
    tf = _box(s, x + 0.1, 3.2, cw - 0.2, 1.0).text_frame; tf.word_wrap = True
    _para(tf, jt, 13, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=3)
    _para(tf, jd, 10.5, WHITE if is_e2e else LIGHTBLUE, align=PP_ALIGN.CENTER)
    if i < 4:
        ar = _box(s, x + cw - 0.02, 2.95, 0.4, 0.6).text_frame
        _para(ar, "›", 26, SLATE, bold=True, align=PP_ALIGN.CENTER, first=True)
    x += cw + 0.21
note = card(s, 0.75, 4.6, 11.83, 1.95, fill=CARDBG, accent=False)
_rect(s, 0.75, 4.6, 0.13, 1.95, TEAL, MSO_SHAPE.ROUNDED_RECTANGLE)
tf2 = _box(s, 1.15, 4.8, 11.1, 1.6).text_frame; tf2.word_wrap = True
_para(tf2, "Se dispara en cada push a develop/main y en cada PR a main.", 14.5, DARK,
      first=True, space_after=7)
_para(tf2, "⚠  El pipeline FALLA si falla cualquier prueba E2E.", 14.5, NAVY, bold=True, space_after=7)
_para(tf2, "Evidencia: el reporte HTML de Playwright se sube como artefacto (14 días) y se "
           "guardan screenshots en cada fallo.", 14.5, DARK)
notes(s, "Asi se integra todo. Ante cada push o PR se disparan cinco jobs encadenados: "
         "builds, unitarias con PostgreSQL real, las E2E de Playwright —que levantan la app "
         "completa— y una notificacion a Discord. Si una E2E falla, el pipeline se marca en "
         "rojo. La evidencia queda como artefacto descargable y con screenshots de fallos.")

# ============================================================
# BLOQUE 5 — DEMOSTRACIÓN
# ============================================================
section("Demostración en vivo", 5, "Matías",
        "El sistema funcionando y las pruebas en acción")

demo_slide(1, "Registro e inicio de sesión", "Matías",
           ["Mostrar el registro con validaciones (email inválido, campo vacío)",
            "Iniciar sesión con un usuario preconfigurado",
            "El JWT se guarda y acompaña a cada petición"],
           "Pantalla de login / registro con usuario demo")
demo_slide(2, "Explorar, buscar y filtrar", "Matías",
           ["Listado con oportunidades ya cargadas (no vacío)",
            "Filtrar por tipo: seleccionar 'TUTORIA'",
            "Buscar por texto: escribir 'matemáticas'",
            "Botones Editar/Eliminar solo en cards propias"],
           "Listado de oportunidades + barra de búsqueda con filtro")
demo_slide(3, "Crear una nueva oportunidad", "Matías",
           ["Completar el formulario con datos realistas",
            "Mostrar la validación de campos obligatorios",
            "Guardar y ver que aparece en el listado"],
           "Formulario de creación con todos los campos")
demo_slide(4, "Funcionalidades nuevas (E2/E3)", "Matías",
           ["Postular y ver la notificación generada",
            "Guardar una oportunidad en favoritos",
            "Gestionar postulantes y dejar feedback",
            "Vista de administración (suspender / bloquear)"],
           "Notificaciones, favoritos y panel de admin")
demo_slide(5, "Las pruebas E2E en acción", "Matías",
           ["Ejecutar 'npx playwright test' en la terminal",
            "Ver los 17 tests corriendo en headless",
            "Abrir el reporte HTML de Playwright",
            "Mostrar un trace de una ejecución"],
           "Terminal con los 17 E2E en verde + reporte Playwright")
demo_slide(6, "El pipeline y la notificación", "Matías",
           ["Mostrar el run en GitHub Actions en verde",
            "Descargar el artefacto del reporte E2E",
            "Mostrar el mensaje de estado en Discord",
            "Plan B: capturas de respaldo por si falla la red"],
           "GitHub Actions (5 jobs en verde) + alerta en Discord")

# ============================================================
# BLOQUE 6 — REFLEXIÓN
# ============================================================
section("Reflexión del equipo", 6, "José",
        "Lo aprendido, las dificultades y las decisiones clave")

s = prs.slides.add_slide(BLANK)
header(s, "Reflexión", "Balance del proyecto", "José")
cols = [
    ("LO APRENDIDO", [
        "El E2E atrapa errores que el unitario no ve",
        "Las variables de entorno bien definidas son críticas en CI",
        "Selectores estables (data-testid) = tests robustos",
        "Probar en CI revela problemas que en local no aparecen"]),
    ("DIFICULTADES", [
        "Levantar app + DB + navegador estable en CI",
        "Sincronizar tipos entre frontend y backend",
        "Evitar tests intermitentes (flaky)",
        "Datos de prueba consistentes entre ejecuciones"]),
    ("DECISIONES CLAVE", [
        "Playwright en vez de Selenium",
        "GitHub Actions en vez de Jenkins",
        "Discord en vez de Slack para alertas",
        "El pipeline falla si falla un E2E"]),
]
x = 0.75
for titulo, items in cols:
    cd = _rect(s, x, 2.0, 3.85, 4.55, WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
    _shadow(cd, blur=11, dist=4, alpha=80)
    head = _rect(s, x, 2.0, 3.85, 0.68, TEAL, MSO_SHAPE.ROUNDED_RECTANGLE)
    head.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _para(head.text_frame, titulo, 12.5, NAVY, bold=True, align=PP_ALIGN.CENTER,
          first=True, spacing=1)
    tc = _box(s, x + 0.28, 2.85, 3.32, 3.6).text_frame; tc.word_wrap = True
    for i, it in enumerate(items):
        p = _para(tc, "•  ", 13, TEAL, bold=True, first=(i == 0), space_after=11)
        r = p.add_run(); r.text = it
        r.font.size = Pt(13); r.font.color.rgb = DARK; r.font.name = 'Calibri'
    x += 4.15
notes(s, "Reflexionando: lo mas valioso fue comprobar que el E2E atrapa errores que el "
         "unitario jamas veria, porque prueba la app completa. La mayor dificultad fue "
         "estabilizar app, base y navegador dentro de CI. Y nuestras decisiones clave "
         "—que defendemos— fueron Playwright sobre Selenium, GitHub Actions sobre Jenkins, "
         "y Discord para las notificaciones.")

s = prs.slides.add_slide(BLANK)
header(s, "Caso real de la Entrega 3", "Un problema que resolvimos", "José")
labels = [("EL SÍNTOMA", "El job de E2E fallaba en CI aunque los tests pasaban en local."),
          ("LA CAUSA", "El health-check del pipeline esperaba respuesta de GET /api, que "
           "devuelve 404 porque no existe esa ruta raíz. La herramienta de espera solo "
           "acepta 2xx, así que expiraba y mataba el job — aunque el backend estaba sano."),
          ("LA SOLUCIÓN", "Apuntar la espera a GET /api/opportunities, un endpoint real que "
           "responde 200. Pipeline en verde y E2E corriendo de punta a punta.")]
y = 2.0
for tit, txt in labels:
    cd = card(s, 0.75, y, 11.83, 1.42, fill=CARDBG, accent=False)
    _rect(s, 0.75, y, 0.13, 1.42, TEAL, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = _box(s, 1.15, y + 0.16, 11.2, 1.15).text_frame; tf.word_wrap = True
    _para(tf, tit, 13, TEALDK, bold=True, first=True, space_after=4, spacing=1.5)
    _para(tf, txt, 14.5, DARK)
    y += 1.6
notes(s, "Un ejemplo concreto, muy del espiritu del curso: el pipeline daba rojo pero en "
         "local todo pasaba. El chequeo de 'backend listo' apuntaba a una ruta que devolvia "
         "404, y la herramienta solo daba por listo un 200. El backend estaba sano; el "
         "chequeo estaba mal. Lo corregimos apuntando a un endpoint real. Moraleja: la "
         "infraestructura de pruebas tambien se prueba y depura.")

# ============================================================
# BLOQUE 7 — PRÓXIMOS PASOS Y CIERRE
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "Cierre", "¿Qué sigue para CampusLink?", "José")
bullets(s, [
    "Ampliar la cobertura E2E a admin, notificaciones y chat",
    "Pruebas E2E multi-navegador (Firefox y WebKit)",
    "Regresión visual de las vistas principales",
    "Despliegue automático tras un pipeline en verde",
    "Métricas de cobertura combinada (unit + E2E) en CI",
], top=2.05, gap=13)
notes(s, "Hacia adelante: ampliar la cobertura E2E a admin, notificaciones y chat, correr "
         "en mas navegadores, sumar regresion visual y automatizar el despliegue cuando el "
         "pipeline queda verde. La base ya esta; ahora es escalarla.")

# Slide de cierre = copia exacta de la portada
s = prs.slides.add_slide(BLANK); DARK_SLIDES.append(s)
s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
_grad_fill(s.background.fill, HEX["NAVY"], HEX["NAVYDK"], angle=60)
_circle(s, 8.7, -2.0, 6.8, HEX["TEAL"], 10)
_circle(s, 10.8, 4.4, 3.8, HEX["TEAL"], 8)
_circle(s, -1.8, 5.0, 4.0, HEX["NAVY2"], 45)
_rect(s, 0, 0, 0.32, 7.5, TEAL)
tf = _box(s, 1.0, 1.95, 11.5, 1.5).text_frame
_para(tf, "CampusLink", 58, WHITE, bold=True, first=True)
tf2 = _box(s, 1.0, 3.5, 11.5, 0.8).text_frame
_para(tf2, "Plataforma de Oportunidades Universitarias", 23, TEAL, first=True)
info = _rect(s, 1.0, 4.75, 8.6, 1.7, NAVYDK, MSO_SHAPE.ROUNDED_RECTANGLE)
_alpha_fill(info, "0F2440", 55); _shadow(info, blur=12, dist=4, alpha=60)
tf3 = _box(s, 1.35, 4.95, 8.0, 1.4).text_frame
_para(tf3, "INF331 – Pruebas de Software   ·   Entrega 3", 15, WHITE, bold=True, first=True, space_after=7)
_para(tf3, "Pruebas E2E/UI con Playwright integradas al pipeline CI/CD", 14, TEAL, space_after=10)
_para(tf3, "José Meza · Matías Barraza · Leonardo Chacón · Aarón Vargas", 13, LIGHTBLUE)
notes(s, "[Habla: José] Cierre. Con esto damos por finalizada la presentacion y abrimos la "
         "ronda de preguntas. Gracias por su atencion.")

# ---------- numeros de pagina ----------
for idx, sl in enumerate(prs.slides, 1):
    footer(sl, idx, dark=(sl in DARK_SLIDES))

out = '/Users/matiasbrrz/Universidad/campuslink/CampusLink_Presentacion_E3.pptx'
prs.save(out)
print(f"Guardado: {out}  ({len(prs.slides._sldIdLst)} slides)")
